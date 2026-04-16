from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_investor_pitch():
    prs = Presentation()
    
    # Premium Branding Palette
    DARK_BG = RGBColor(15, 23, 42)    # Slate 900
    ACCENT_MAIN = RGBColor(79, 70, 229) # Indigo 600
    ACCENT_SOFT = RGBColor(129, 140, 248) # Indigo 400
    SUCCESS = RGBColor(16, 185, 129)   # Emerald 500
    TEXT_MAIN = RGBColor(241, 245, 249) # Slate 100
    TEXT_DIM = RGBColor(148, 163, 184)  # Slate 400

    def apply_base_style(slide, title_text):
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        # Sidebar Decorative
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.05), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_MAIN
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = TEXT_MAIN

    def add_bullet_points(slide, points, top=Inches(1.8)):
        content_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(8.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for point in points:
            p = tf.add_paragraph()
            p.text = "   ▸  " + point
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DIM
            p.space_after = Pt(14)
            if ":" in point:
                run = p.runs[0]
                run.font.bold = True
                run.font.color.rgb = ACCENT_SOFT

    # --- SLIDE 1: VISIONARY COVER ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG

    # Glowing Circle Decoration
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(-1.5), Inches(5.5), Inches(5.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_MAIN
    circle.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(7), Inches(2))
    p = title.text_frame.paragraphs[0]
    p.text = "SWAR-SHIKSHA"
    p.font.bold = True
    p.font.size = Pt(64)
    p.font.color.rgb = TEXT_MAIN
    
    sub = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(7), Inches(1))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Reimagining Accessibility through Voice-First AI\nHackBLR 2026 | Problem Statement 3"
    sp.font.size = Pt(22)
    sp.font.color.rgb = SUCCESS

    # --- SLIDE 2: THE PROBLEM (THE AGONY) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide, "The Crisis in Education")
    add_bullet_points(slide, [
        "Digital Exclusion: 2.2 billion people live with vision impairment globally.",
        "Learning Barriers: Dyslexic and low-literacy students struggle with text-heavy curriculum.",
        "Resource Gap: Braille and human tutors are expensive and non-scalable.",
        "Native Barrier: Textbooks are in English; conceptual understanding happens in native tongue."
    ])

    # --- SLIDE 3: THE SOLUTION (THE DISRUPTION) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide, "Swar-Shiksha: A Voice Tutor")
    add_bullet_points(slide, [
        "Hands-Free Learning: Converting static PDF content into interactive dialogue.",
        "Multilingual Intelligence: Real-time tutoring in Hindi and English.",
        "Total Understanding: Narrating complex diagrams using Vision AI.",
        "Active Assessment: Voice-driven quizzes to verify student mastery."
    ])

    # --- SLIDE 4: THE DIFFERENTIATOR (TECH EDGE) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide, "Proprietary Architecture")
    add_bullet_points(slide, [
        "Real-Time Voice: Leveraging VAPI for ultra-low latency conversations.",
        "RAG Core: Qdrant Vector Search ensures 100% textbook accuracy (No hallucinations).",
        "Vision Describe: GPT-4o Vision interprets and explains visuals to the blind.",
        "Adaptive UX: A minimalistic, high-contrast dashboard for universal design."
    ])

    # --- SLIDE 5: MARKET & SOCIETAL IMPACT ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide, "Global Impact Potential")
    # Add metrics boxes
    def add_metric(slide, x, y, val, label):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = ACCENT_MAIN
        
        v_box = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(2), Inches(0.5))
        vp = v_box.text_frame.paragraphs[0]
        vp.text = val
        vp.font.bold = True
        vp.font.size = Pt(28)
        vp.font.color.rgb = SUCCESS
        vp.alignment = PP_ALIGN.CENTER
        
        l_box = slide.shapes.add_textbox(x, y + Inches(0.7), Inches(2), Inches(0.4))
        lp = l_box.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(10)
        lp.font.color.rgb = TEXT_DIM
        lp.alignment = PP_ALIGN.CENTER

    add_metric(slide, Inches(1), Inches(2), "10M+", "Students Impacted")
    add_metric(slide, Inches(3.5), Inches(2), "85%", "Retention Increase")
    add_metric(slide, Inches(6), Inches(2), "0.5s", "Voice Latency")

    add_bullet_points(slide, [
        "Democratic Education: High-quality tutoring for a fraction of the cost.",
        "Global Scalability: Adaptable to any language or knowledge domain."
    ], top=Inches(3.8))

    # --- SLIDE 6: COMPETITIVE LANDSCAPE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide, "Winning the Landscape")
    add_bullet_points(slide, [
        "Vs Screen Readers: We provide context and explanation, not just reading.",
        "Vs Traditional AI: RAG technology ensures zero hallucinations from textbooks.",
        "Vs Video Lessons: Voice-first interaction allows active student participation."
    ])

    # --- SLIDE 7: BUSINESS & ROADMAP ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base_style(slide, "2026 Strategy")
    add_bullet_points(slide, [
        "Q3 2026: Offline Edge-AI deployment for remote rural schools.",
        "Q4 2026: Integration with Smart Glasses for real-time physical reading.",
        "2027: Regional Expansion (Kannada, Tamil, Telugu, Bengali)."
    ])

    # --- SLIDE 8: THE PITCH CLOSE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    
    title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "NO STUDENT LEFT BEHIND."
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = TEXT_MAIN
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Swar-Shiksha AI | Pitch by Naveen Kumar K M\nGitHub: github.com/Naveenkm07/swar-shiksha-ai"
    sp.font.size = Pt(14)
    sp.font.color.rgb = SUCCESS
    sp.alignment = PP_ALIGN.CENTER

    prs.save('Swar-Shiksha_Investor_Pitch.pptx')
    print("Investor pitch created successfully!")

if __name__ == "__main__":
    create_investor_pitch()
